"""
Core logic for Instrumentl Grant Matcher.
Used by both the desktop app (grant_matcher.py) and the Streamlit web app.
"""

import os
import json
import time
import re
from collections import Counter
import math
from pathlib import Path

# ==============================================================================
# CONFIGURATION
# ==============================================================================

CONFIG_FILE = "config.json"
DEFAULT_CONFIG = {
    "api_key_id": "",
    "api_private_key": "",
    "chunk_size": 500,
    "min_match_score": 0.1,
    "top_matches": 100,
    "last_export_dir": ""
}


def load_config():
    """Load configuration from file."""
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)
                for key, value in DEFAULT_CONFIG.items():
                    if key not in config:
                        config[key] = value
                return config
        except:
            pass
    return DEFAULT_CONFIG.copy()


def save_config(config):
    """Save configuration to file."""
    with open(CONFIG_FILE, 'w') as f:
        json.dump(config, f, indent=2)


# ==============================================================================
# DOCUMENT PROCESSING
# ==============================================================================

class DocumentProcessor:
    """Handles document reading and text extraction."""

    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.csv', '.txt', '.md'}

    @staticmethod
    def extract_text(filepath):
        """Extract text from various document types."""
        ext = Path(filepath).suffix.lower()
        try:
            if ext == '.txt' or ext == '.md':
                return DocumentProcessor._read_text_file(filepath)
            elif ext == '.pdf':
                return DocumentProcessor._read_pdf(filepath)
            elif ext in ('.docx', '.doc'):
                return DocumentProcessor._read_word(filepath)
            elif ext in ('.xlsx', '.xls'):
                return DocumentProcessor._read_excel(filepath)
            elif ext in ('.pptx', '.ppt'):
                return DocumentProcessor._read_powerpoint(filepath)
            elif ext == '.csv':
                return DocumentProcessor._read_csv(filepath)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            raise Exception(f"Error reading {filepath}: {str(e)}")

    @staticmethod
    def _read_text_file(filepath):
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise ValueError("Could not decode text file")

    @staticmethod
    def _read_pdf(filepath):
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            return '\n\n'.join(text_parts)
        except ImportError:
            from pypdf import PdfReader
            reader = PdfReader(filepath)
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            return '\n\n'.join(text_parts)

    @staticmethod
    def _read_word(filepath):
        ext = Path(filepath).suffix.lower()
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(filepath)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return '\n\n'.join(paragraphs)
            except ImportError:
                import subprocess
                result = subprocess.run(['pandoc', filepath, '-t', 'plain'], capture_output=True, text=True)
                if result.returncode == 0:
                    return result.stdout
                raise Exception("python-docx not installed and pandoc failed")
        else:
            import subprocess
            result = subprocess.run(['pandoc', filepath, '-t', 'plain'], capture_output=True, text=True)
            if result.returncode == 0:
                return result.stdout
            raise Exception("Cannot read .doc files - install pandoc")

    @staticmethod
    def _read_excel(filepath):
        try:
            import pandas as pd
            xlsx = pd.ExcelFile(filepath)
            text_parts = []
            for sheet_name in xlsx.sheet_names:
                df = pd.read_excel(xlsx, sheet_name=sheet_name)
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(df.to_string())
            return '\n\n'.join(text_parts)
        except ImportError:
            raise Exception("pandas not installed - cannot read Excel files")

    @staticmethod
    def _read_powerpoint(filepath):
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"Slide {slide_num}:"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if len(slide_text) > 1:
                    text_parts.append('\n'.join(slide_text))
            return '\n\n'.join(text_parts)
        except ImportError:
            raise Exception("python-pptx not installed - cannot read PowerPoint files")

    @staticmethod
    def _read_csv(filepath):
        try:
            import pandas as pd
            df = pd.read_csv(filepath)
            return df.to_string()
        except ImportError:
            import csv
            with open(filepath, 'r', newline='', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = [', '.join(row) for row in reader]
            return '\n'.join(rows)


class TextChunker:
    @staticmethod
    def chunk_text(text, chunk_size=500, overlap=50):
        text = re.sub(r'\s+', ' ', text).strip()
        words = text.split()
        if len(words) <= chunk_size:
            return [text] if text else []
        chunks = []
        start = 0
        while start < len(words):
            end = start + chunk_size
            if end < len(words):
                search_start = start + int(chunk_size * 0.8)
                best_break = end
                for i in range(search_start, end):
                    if words[i].endswith(('.', '!', '?', ':', ';')):
                        best_break = i + 1
                end = best_break
            chunk = ' '.join(words[start:end])
            chunks.append(chunk)
            start = end - overlap if end - overlap > start else end
        return chunks


# ==============================================================================
# TF-IDF MATCHING
# ==============================================================================

class TFIDFMatcher:
    def __init__(self):
        self.documents = []
        self.doc_metadata = []
        self.vocabulary = {}
        self.idf = {}
        self.doc_vectors = []

    def add_documents(self, documents, metadata_list=None):
        if metadata_list is None:
            metadata_list = [{}] * len(documents)
        for doc, meta in zip(documents, metadata_list):
            self.documents.append(doc)
            self.doc_metadata.append(meta)

    def build_index(self):
        if not self.documents:
            return
        tokenized_docs = [self._tokenize(doc) for doc in self.documents]
        all_terms = set()
        for tokens in tokenized_docs:
            all_terms.update(tokens)
        self.vocabulary = {term: idx for idx, term in enumerate(sorted(all_terms))}
        doc_count = len(self.documents)
        doc_freq = Counter()
        for tokens in tokenized_docs:
            unique_tokens = set(tokens)
            for token in unique_tokens:
                doc_freq[token] += 1
        self.idf = {}
        for term, freq in doc_freq.items():
            self.idf[term] = math.log((doc_count + 1) / (freq + 1)) + 1
        self.doc_vectors = []
        for tokens in tokenized_docs:
            vector = self._calculate_tfidf_vector(tokens)
            self.doc_vectors.append(vector)

    def find_matches(self, query_text, top_k=10, min_score=0.0):
        if not self.doc_vectors:
            return []
        query_tokens = self._tokenize(query_text)
        query_vector = self._calculate_tfidf_vector(query_tokens)
        scores = []
        for idx, doc_vector in enumerate(self.doc_vectors):
            score = self._cosine_similarity(query_vector, doc_vector)
            if score >= min_score:
                scores.append((idx, score))
        scores.sort(key=lambda x: x[1], reverse=True)
        results = []
        for idx, score in scores[:top_k]:
            results.append({
                'document': self.documents[idx],
                'metadata': self.doc_metadata[idx],
                'score': score
            })
        return results

    def _tokenize(self, text):
        text = text.lower()
        text = re.sub(r'[^a-z0-9\s]', ' ', text)
        words = text.split()
        stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
                      'of', 'with', 'by', 'from', 'is', 'are', 'was', 'were', 'be', 'been',
                      'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would',
                      'could', 'should', 'may', 'might', 'must', 'shall', 'can', 'this',
                      'that', 'these', 'those', 'it', 'its', 'as', 'if', 'then', 'than',
                      'so', 'such', 'no', 'not', 'only', 'same', 'too', 'very', 'just',
                      'also', 'into', 'over', 'after', 'before', 'above', 'below', 'up',
                      'down', 'out', 'off', 'about', 'each', 'all', 'any', 'both', 'few',
                      'more', 'most', 'other', 'some', 'their', 'them', 'they', 'we', 'us',
                      'our', 'your', 'you', 'he', 'she', 'him', 'her', 'his', 'hers'}
        tokens = [w for w in words if len(w) > 2 and w not in stop_words]
        return tokens

    def _calculate_tfidf_vector(self, tokens):
        tf = Counter(tokens)
        total_terms = len(tokens) if tokens else 1
        vector = {}
        for term, count in tf.items():
            if term in self.vocabulary:
                tf_value = count / total_terms
                idf_value = self.idf.get(term, 1.0)
                vector[term] = tf_value * idf_value
        return vector

    def _cosine_similarity(self, vec1, vec2):
        common_terms = set(vec1.keys()) & set(vec2.keys())
        if not common_terms:
            return 0.0
        dot_product = sum(vec1[term] * vec2[term] for term in common_terms)
        mag1 = math.sqrt(sum(v ** 2 for v in vec1.values()))
        mag2 = math.sqrt(sum(v ** 2 for v in vec2.values()))
        if mag1 == 0 or mag2 == 0:
            return 0.0
        return dot_product / (mag1 * mag2)


# ==============================================================================
# INSTRUMENTL API CLIENT
# ==============================================================================

class InstrumentlAPI:
    """Client for the Instrumentl API."""

    BASE_URL = "https://api.instrumentl.com"

    def __init__(self, api_key_id, api_private_key):
        self.api_key_id = api_key_id
        self.api_private_key = api_private_key
        self._session = None
        self._init_session()

    def _init_session(self):
        try:
            import requests
            self._session = requests.Session()
            self._session.auth = (self.api_key_id, self.api_private_key)
            self._session.headers.update({
                'Accept': 'application/json',
                'Content-Type': 'application/json',
                'User-Agent': 'Mozilla/5.0 InstrumentlGrantMatcher/2.0'
            })
            self._use_requests = True
        except ImportError:
            self._use_requests = False

    def _make_request(self, endpoint, params=None):
        url = f"{self.BASE_URL}{endpoint}"
        if self._use_requests:
            return self._make_request_with_requests(url, params)
        else:
            return self._make_request_with_urllib(url, params)

    def _make_request_with_requests(self, url, params=None):
        import requests
        try:
            response = self._session.get(url, params=params, timeout=30)
            if response.status_code == 403:
                raise Exception("API Error 403: Access denied.")
            elif response.status_code == 402:
                raise Exception("API Error 402: No active API subscription.")
            elif response.status_code == 401:
                raise Exception("API Error 401: Invalid API credentials.")
            elif response.status_code == 404:
                return None
            response.raise_for_status()
            return response.json()
        except requests.exceptions.HTTPError as e:
            error_text = e.response.text if e.response else str(e)
            raise Exception(f"API Error: {error_text}")
        except requests.exceptions.ConnectionError:
            raise Exception("Connection Error: Could not connect to Instrumentl API.")
        except requests.exceptions.Timeout:
            raise Exception("Timeout Error: Request timed out.")

    def _make_request_with_urllib(self, url, params=None):
        import urllib.request
        import urllib.parse
        import base64
        if params:
            url += "?" + urllib.parse.urlencode(params)
        request = urllib.request.Request(url)
        auth_string = f"{self.api_key_id}:{self.api_private_key}"
        auth_b64 = base64.b64encode(auth_string.encode('utf-8')).decode('utf-8')
        request.add_header('Authorization', f'Basic {auth_b64}')
        request.add_header('Accept', 'application/json')
        try:
            with urllib.request.urlopen(request, timeout=30) as response:
                return json.loads(response.read().decode('utf-8'))
        except urllib.error.HTTPError as e:
            if e.code == 404:
                return None
            raise Exception(f"API Error {e.code}")
        except urllib.error.URLError as e:
            raise Exception(f"Connection Error: {str(e.reason)}")

    def get_account(self):
        return self._make_request("/v1/accounts/current")

    def get_projects(self, page_size=50, cursor=None):
        params = {"page_size": page_size}
        if cursor:
            params["cursor"] = cursor
        return self._make_request("/v1/projects", params)

    def get_all_projects(self, callback=None):
        all_projects = []
        cursor = None
        page = 1
        while True:
            if callback:
                callback(f"Fetching projects page {page}...")
            result = self.get_projects(page_size=50, cursor=cursor)
            if not result:
                break
            projects = result.get('projects', [])
            for p in projects:
                if _is_active_project(p):
                    all_projects.append(p)
            meta = result.get('meta', {})
            if not meta.get('has_more', False):
                break
            cursor = meta.get('cursor')
            page += 1
            time.sleep(0.25)
        return all_projects

    def get_grants(self, page_size=50, cursor=None, is_saved=None, funder_id=None):
        params = {"page_size": page_size}
        if cursor:
            params["cursor"] = cursor
        if is_saved is not None:
            params["is_saved"] = str(is_saved).lower()
        if funder_id:
            params["funder_id"] = funder_id
        return self._make_request("/v1/grants", params)

    def get_grant(self, grant_id):
        return self._make_request(f"/v1/grants/{grant_id}")

    def get_saved_grants(self, page_size=50, cursor=None, project_id=None):
        params = {"page_size": page_size}
        if cursor:
            params["cursor"] = cursor
        if project_id:
            params["project_id"] = project_id
        return self._make_request("/v1/saved_grants", params)

    def get_all_grants(self, callback=None):
        all_grants = []
        cursor = None
        page = 1
        while True:
            if callback:
                callback(f"Fetching grants page {page}...")
            result = self.get_grants(page_size=50, cursor=cursor)
            if not result:
                break
            grants = result.get('grants', [])
            all_grants.extend(grants)
            meta = result.get('meta', {})
            if not meta.get('has_more', False):
                break
            cursor = meta.get('cursor')
            page += 1
            time.sleep(0.25)
        return all_grants

    def get_all_saved_grants(self, project_id=None, callback=None):
        all_saved = []
        cursor = None
        page = 1
        while True:
            if callback:
                callback(f"Fetching saved grants page {page}...")
            result = self.get_saved_grants(page_size=50, cursor=cursor, project_id=project_id)
            if not result:
                break
            saved = result.get('saved_grants', [])
            all_saved.extend(saved)
            meta = result.get('meta', {})
            if not meta.get('has_more', False):
                break
            cursor = meta.get('cursor')
            page += 1
            time.sleep(0.25)
        return all_saved


# ==============================================================================
# PROJECT FILTER
# ==============================================================================

def _is_active_project(project):
    """Return True if a project should be considered active.

    The Instrumentl API may include archived or deleted projects in the
    /v1/projects response.  We inspect several common status fields and
    exclude any project that is clearly inactive.  If none of these fields
    are present on the object we assume the project is active (fail-open).
    """
    # status field: keep only if value is 'active' (case-insensitive)
    status = project.get('status')
    if status is not None:
        return str(status).lower() == 'active'

    # archived_at timestamp: non-null means archived
    archived_at = project.get('archived_at')
    if archived_at is not None:
        return False

    # archived boolean flag
    archived = project.get('archived')
    if archived:
        return False

    # is_active boolean flag
    is_active = project.get('is_active')
    if is_active is not None:
        return bool(is_active)

    # No status field found — assume active
    return True


# ==============================================================================
# LOCATION FILTER
# ==============================================================================

def grant_matches_location(grant, location_filter):
    """Check if grant matches the location filter."""
    categories = grant.get('categories', {})
    if not isinstance(categories, dict):
        return True
    geographic_area = categories.get('geographic_area_category', [])
    if not geographic_area:
        return True
    locations = [loc.lower() if isinstance(loc, str) else '' for loc in geographic_area]
    if location_filter == "indiana":
        return any('indiana' in loc or 'in' == loc for loc in locations)
    elif location_filter == "usa":
        nationwide_terms = ['national', 'nationwide', 'united states', 'usa', 'all states', 'u.s.']
        return any(any(term in loc for term in nationwide_terms) for loc in locations)
    elif location_filter == "indiana_usa":
        has_indiana = any('indiana' in loc or 'in' == loc for loc in locations)
        nationwide_terms = ['national', 'nationwide', 'united states', 'usa', 'all states', 'u.s.']
        has_nationwide = any(any(term in loc for term in nationwide_terms) for loc in locations)
        return has_indiana or has_nationwide
    return True


# ==============================================================================
# RESULTS HELPER
# ==============================================================================

def build_results_dataframe(match_results):
    """Convert match results list into a pandas DataFrame for display/export."""
    import pandas as pd
    rows = []
    for rank, result in enumerate(match_results, 1):
        grant = result['metadata']
        funder = grant.get('funder', '')
        funder_name = funder.get('name', '') if isinstance(funder, dict) else str(funder)
        cycles = grant.get('funding_cycles', [])
        funding_cycle = cycles[0].get('interval', '') if cycles else ''
        slug = grant.get('slug', '')
        grant_url = f"https://www.instrumentl.com/grants/{slug}" if slug else ''
        rows.append({
            'Rank': rank,
            'Score': round(result['score'], 4),
            'Grant Name': grant.get('name', ''),
            'Grant ID': grant.get('id', ''),
            'Funder': funder_name,
            'Next Deadline': grant.get('next_deadline_date', ''),
            'Status': grant.get('status', ''),
            'Is Custom': grant.get('is_custom', False),
            'Rolling': grant.get('rolling', False),
            'Funding Cycle': funding_cycle,
            'Grant URL': grant_url,
            'Description': (grant.get('overview', '') or '')[:500],
        })
    return pd.DataFrame(rows)

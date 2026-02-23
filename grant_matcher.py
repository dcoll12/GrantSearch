#!/usr/bin/env python3
"""
Instrumentl Grant Matcher Application v2.0
==========================================
A desktop application that matches your organization's documents
against grant opportunities from Instrumentl.

No coding required - just click Run!
"""

import os
import sys
import json
import time
import threading
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
from pathlib import Path
from datetime import datetime
import re
from collections import Counter
import math
import random

# ==============================================================================
# CONFIGURATION
# ==============================================================================

CONFIG_FILE = "config.json"
DEFAULT_CONFIG = {
    "api_key_id": "",
    "api_private_key": "",
    "chunk_size": 500,
    "min_match_score": 0.1,
    "top_matches": 100,
    "last_export_dir": "",
    "max_retries": 3,
    "retry_base_delay": 1.0,
    "retry_max_delay": 30.0
}


def load_config():
    """Load configuration from file."""
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)
                for key, value in DEFAULT_CONFIG.items():
                    if key not in config:
                        config[key] = value
                return config
        except:
            pass
    return DEFAULT_CONFIG.copy()


def save_config(config):
    """Save configuration to file."""
    with open(CONFIG_FILE, 'w') as f:
        json.dump(config, f, indent=2)


# ==============================================================================
# DOCUMENT PROCESSING
# ==============================================================================

class DocumentProcessor:
    """Handles document reading and text extraction."""

    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.csv', '.txt', '.md'}

    @staticmethod
    def extract_text(filepath):
        """Extract text from various document types."""
        ext = Path(filepath).suffix.lower()

        try:
            if ext == '.txt' or ext == '.md':
                return DocumentProcessor._read_text_file(filepath)
            elif ext == '.pdf':
                return DocumentProcessor._read_pdf(filepath)
            elif ext in ('.docx', '.doc'):
                return DocumentProcessor._read_word(filepath)
            elif ext in ('.xlsx', '.xls'):
                return DocumentProcessor._read_excel(filepath)
            elif ext in ('.pptx', '.ppt'):
                return DocumentProcessor._read_powerpoint(filepath)
            elif ext == '.csv':
                return DocumentProcessor._read_csv(filepath)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            raise Exception(f"Error reading {filepath}: {str(e)}")

    @staticmethod
    def _read_text_file(filepath):
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise ValueError("Could not decode text file")

    @staticmethod
    def _read_pdf(filepath):
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            return '\n\n'.join(text_parts)
        except ImportError:
            from pypdf import PdfReader
            reader = PdfReader(filepath)
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            return '\n\n'.join(text_parts)

    @staticmethod
    def _read_word(filepath):
        ext = Path(filepath).suffix.lower()
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(filepath)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                return '\n\n'.join(paragraphs)
            except ImportError:
                import subprocess
                result = subprocess.run(['pandoc', filepath, '-t', 'plain'], capture_output=True, text=True)
                if result.returncode == 0:
                    return result.stdout
                raise Exception("python-docx not installed and pandoc failed")
        else:
            import subprocess
            result = subprocess.run(['pandoc', filepath, '-t', 'plain'], capture_output=True, text=True)
            if result.returncode == 0:
                return result.stdout
            raise Exception("Cannot read .doc files - install pandoc")

    @staticmethod
    def _read_excel(filepath):
        try:
            import pandas as pd
            xlsx = pd.ExcelFile(filepath)
            text_parts = []
            for sheet_name in xlsx.sheet_names:
                df = pd.read_excel(xlsx, sheet_name=sheet_name)
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(df.to_string())
            return '\n\n'.join(text_parts)
        except ImportError:
            raise Exception("pandas not installed - cannot read Excel files")

    @staticmethod
    def _read_powerpoint(filepath):
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"Slide {slide_num}:"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if len(slide_text) > 1:
                    text_parts.append('\n'.join(slide_text))
            return '\n\n'.join(text_parts)
        except ImportError:
            raise Exception("python-pptx not installed - cannot read PowerPoint files")

    @staticmethod
    def _read_csv(filepath):
        try:
            import pandas as pd
            df = pd.read_csv(filepath)
            return df.to_string()
        except ImportError:
            import csv
            with open(filepath, 'r', newline='', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = [', '.join(row) for row in reader]
            return '\n'.join(rows)


class TextChunker:
    @staticmethod
    def chunk_text(text, chunk_size=500, overlap=50):
        text = re.sub(r'\s+', ' ', text).strip()
        words = text.split()
        if len(words) <= chunk_size:
            return [text] if text else []
        chunks = []
        start = 0
        while start < len(words):
            end = start + chunk_size
            if end < len(words):
                search_start = start + int(chunk_size * 0.8)
                best_break = end
                for i in range(search_start, end):
                    if words[i].endswith(('.', '!', '?', ':', ';')):
                        best_break = i + 1
                end = best_break
            chunk = ' '.join(words[start:end])
            chunks.append(chunk)
            start = end - overlap if end - overlap > start else end
        return chunks


# ==============================================================================
# TF-IDF MATCHING
# ==============================================================================

class TFIDFMatcher:
    def __init__(self):
        self.documents = []
        self.doc_metadata = []
        self.vocabulary = {}
        self.idf = {}
        self.doc_vectors = []

    def add_documents(self, documents, metadata_list=None):
        if metadata_list is None:
            metadata_list = [{}] * len(documents)
        for doc, meta in zip(documents, metadata_list):
            self.documents.append(doc)
            self.doc_metadata.append(meta)

    def build_index(self):
        if not self.documents:
            return
        tokenized_docs = [self._tokenize(doc) for doc in self.documents]
        all_terms = set()
        for tokens in tokenized_docs:
            all_terms.update(tokens)
        self.vocabulary = {term: idx for idx, term in enumerate(sorted(all_terms))}
        doc_count = len(self.documents)
        doc_freq = Counter()
        for tokens in tokenized_docs:
            unique_tokens = set(tokens)
            for token in unique_tokens:
                doc_freq[token] += 1
        self.idf = {}
        for term, freq in doc_freq.items():
            self.idf[term] = math.log((doc_count + 1) / (freq + 1)) + 1
        self.doc_vectors = []
        for tokens in tokenized_docs:
            vector = self._calculate_tfidf_vector(tokens)
            self.doc_vectors.append(vector)

    def find_matches(self, query_text, top_k=10, min_score=0.0):
        if not self.doc_vectors:
            return []
        query_tokens = self._tokenize(query_text)
        query_vector = self._calculate_tfidf_vector(query_tokens)
        scores = []
        for idx, doc_vector in enumerate(self.doc_vectors):
            score = self._cosine_similarity(query_vector, doc_vector)
            if score >= min_score:
                scores.append((idx, score))
        scores.sort(key=lambda x: x[1], reverse=True)
        results = []
        for idx, score in scores[:top_k]:
            results.append({
                'document': self.documents[idx],
                'metadata': self.doc_metadata[idx],
                'score': score
            })
        return results

    def _tokenize(self, text):
        text = text.lower()
        text = re.sub(r'[^a-z0-9\s]', ' ', text)
        words = text.split()
        stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
                      'of', 'with', 'by', 'from', 'is', 'are', 'was', 'were', 'be', 'been',
                      'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would',
                      'could', 'should', 'may', 'might', 'must', 'shall', 'can', 'this',
                      'that', 'these', 'those', 'it', 'its', 'as', 'if', 'then', 'than',
                      'so', 'such', 'no', 'not', 'only', 'same', 'too', 'very', 'just',
                      'also', 'into', 'over', 'after', 'before', 'above', 'below', 'up',
                      'down', 'out', 'off', 'about', 'each', 'all', 'any', 'both', 'few',
                      'more', 'most', 'other', 'some', 'their', 'them', 'they', 'we', 'us',
                      'our', 'your', 'you', 'he', 'she', 'him', 'her', 'his', 'hers'}
        tokens = [w for w in words if len(w) > 2 and w not in stop_words]
        return tokens

    def _calculate_tfidf_vector(self, tokens):
        tf = Counter(tokens)
        total_terms = len(tokens) if tokens else 1
        vector = {}
        for term, count in tf.items():
            if term in self.vocabulary:
                tf_value = count / total_terms
                idf_value = self.idf.get(term, 1.0)
                vector[term] = tf_value * idf_value
        return vector

    def _cosine_similarity(self, vec1, vec2):
        common_terms = set(vec1.keys()) & set(vec2.keys())
        if not common_terms:
            return 0.0
        dot_product = sum(vec1[term] * vec2[term] for term in common_terms)
        mag1 = math.sqrt(sum(v ** 2 for v in vec1.values()))
        mag2 = math.sqrt(sum(v ** 2 for v in vec2.values()))
        if mag1 == 0 or mag2 == 0:
            return 0.0
        return dot_product / (mag1 * mag2)


# ==============================================================================
# PROJECT FILTER
# ==============================================================================

def _is_active_project(project):
    """Return True if a project should be considered active.

    The Instrumentl API may include archived or deleted projects in the
    /v1/projects response.  We inspect several common status fields and
    exclude any project that is clearly inactive.  If none of these fields
    are present on the object we assume the project is active (fail-open).
    """
    # status field: keep only if value is 'active' (case-insensitive)
    status = project.get('status')
    if status is not None:
        return str(status).lower() == 'active'

    # archived_at timestamp: non-null means archived
    archived_at = project.get('archived_at')
    if archived_at is not None:
        return False

    # archived boolean flag
    archived = project.get('archived')
    if archived:
        return False

    # is_active boolean flag
    is_active = project.get('is_active')
    if is_active is not None:
        return bool(is_active)

    # No status field found — assume active
    return True


# ==============================================================================
# INSTRUMENTL API CLIENT
# ==============================================================================

class RetryableAPIError(Exception):
    """Raised for API errors that are safe to retry (429, 5xx, timeouts, connection errors)."""

    def __init__(self, message, retry_after=None):
        super().__init__(message)
        self.retry_after = retry_after


class InstrumentlAPI:
    """Client for the Instrumentl API."""

    BASE_URL = "https://api.instrumentl.com"

    def __init__(self, api_key_id, api_private_key, max_retries=3, retry_base_delay=1.0, retry_max_delay=30.0):
        self.api_key_id = api_key_id
        self.api_private_key = api_private_key
        self.max_retries = max_retries
        self.retry_base_delay = retry_base_delay
        self.retry_max_delay = retry_max_delay
        self._session = None
        self._init_session()

    def _init_session(self):
        try:
            import requests
            self._session = requests.Session()
            self._session.auth = (self.api_key_id, self.api_private_key)
            self._session.headers.update({
                'Accept': 'application/json',
                'Content-Type': 'application/json',
                'User-Agent': 'Mozilla/5.0 InstrumentlGrantMatcher/2.0'
            })
            self._use_requests = True
        except ImportError:
            self._use_requests = False

    def _make_request(self, endpoint, params=None):
        url = f"{self.BASE_URL}{endpoint}"
        last_error = None
        for attempt in range(self.max_retries + 1):
            try:
                if self._use_requests:
                    return self._make_request_with_requests(url, params)
                else:
                    return self._make_request_with_urllib(url, params)
            except RetryableAPIError as e:
                last_error = e
                if attempt >= self.max_retries:
                    break
                if e.retry_after is not None:
                    delay = min(e.retry_after, self.retry_max_delay)
                else:
                    delay = min(
                        self.retry_base_delay * (2 ** attempt) + random.uniform(0, 1),
                        self.retry_max_delay
                    )
                print(f"Retry {attempt + 1}/{self.max_retries} after {delay:.1f}s: {e}")
                time.sleep(delay)
        raise Exception(f"{last_error} (after {self.max_retries} retries)")

    def _make_request_with_requests(self, url, params=None):
        import requests
        try:
            response = self._session.get(url, params=params, timeout=30)
            if response.status_code == 429:
                retry_after = response.headers.get('Retry-After')
                retry_delay = float(retry_after) if retry_after else None
                raise RetryableAPIError(
                    f"API Error 429: Rate limited.",
                    retry_after=retry_delay
                )
            elif response.status_code >= 500:
                raise RetryableAPIError(
                    f"API Error {response.status_code}: Server error."
                )
            elif response.status_code == 403:
                raise Exception("API Error 403: Access denied.")
            elif response.status_code == 402:
                raise Exception("API Error 402: No active API subscription.")
            elif response.status_code == 401:
                raise Exception("API Error 401: Invalid API credentials.")
            elif response.status_code == 404:
                return None
            response.raise_for_status()
            return response.json()
        except requests.exceptions.HTTPError as e:
            error_text = e.response.text if e.response else str(e)
            raise Exception(f"API Error: {error_text}")
        except requests.exceptions.ConnectionError:
            raise RetryableAPIError("Connection Error: Could not connect to Instrumentl API.")
        except requests.exceptions.Timeout:
            raise RetryableAPIError("Timeout Error: Request timed out.")

    def _make_request_with_urllib(self, url, params=None):
        import urllib.request
        import urllib.error
        import urllib.parse
        import base64
        if params:
            url += "?" + urllib.parse.urlencode(params)
        request = urllib.request.Request(url)
        auth_string = f"{self.api_key_id}:{self.api_private_key}"
        auth_b64 = base64.b64encode(auth_string.encode('utf-8')).decode('utf-8')
        request.add_header('Authorization', f'Basic {auth_b64}')
        request.add_header('Accept', 'application/json')
        try:
            with urllib.request.urlopen(request, timeout=30) as response:
                return json.loads(response.read().decode('utf-8'))
        except urllib.error.HTTPError as e:
            if e.code == 404:
                return None
            if e.code == 429:
                retry_after = e.headers.get('Retry-After') if hasattr(e, 'headers') else None
                retry_delay = float(retry_after) if retry_after else None
                raise RetryableAPIError(
                    f"API Error 429: Rate limited.",
                    retry_after=retry_delay
                )
            if e.code >= 500:
                raise RetryableAPIError(f"API Error {e.code}: Server error.")
            raise Exception(f"API Error {e.code}")
        except urllib.error.URLError as e:
            raise RetryableAPIError(f"Connection Error: {str(e.reason)}")

    def get_account(self):
        return self._make_request("/v1/accounts/current")

    def get_projects(self, page_size=50, cursor=None):
        params = {"page_size": page_size}
        if cursor:
            params["cursor"] = cursor
        return self._make_request("/v1/projects", params)

    def get_all_projects(self, callback=None):
        all_projects = []
        seen_ids = set()
        cursor = None
        prev_cursor = None
        page = 1
        while True:
            if callback:
                callback(f"Fetching projects page {page}...")
            result = self.get_projects(page_size=50, cursor=cursor)
            if not result:
                break
            projects = result.get('projects', [])
            print(f"[DEBUG] Page {page}: got {len(projects)} projects, keys={list(result.keys())}")
            if page == 1 and projects:
                print(f"[DEBUG] First project keys: {list(projects[0].keys())}")
            for p in projects:
                pid = p.get('id')
                if pid not in seen_ids and _is_active_project(p):
                    seen_ids.add(pid)
                    all_projects.append(p)
            meta = result.get('meta', {})
            print(f"[DEBUG] Meta: {meta}")
            if not meta.get('has_more', False):
                break
            new_cursor = meta.get('cursor')
            if new_cursor == prev_cursor:
                print(f"[DEBUG] Cursor unchanged, stopping pagination")
                break
            prev_cursor = cursor
            cursor = new_cursor
            page += 1
            time.sleep(0.25)
        print(f"[DEBUG] Total unique projects: {len(all_projects)}")
        return all_projects

    def get_grants(self, page_size=50, cursor=None, is_saved=None, funder_id=None):
        params = {"page_size": page_size}
        if cursor:
            params["cursor"] = cursor
        if is_saved is not None:
            params["is_saved"] = str(is_saved).lower()
        if funder_id:
            params["funder_id"] = funder_id
        return self._make_request("/v1/grants", params)

    def get_grants_first_page(self, project_id=None):
        """Fetch only the first page of grants (up to 50).

        When a project_id is provided it is passed as a query parameter so the
        API can return grants matched to that specific project (the same set
        shown on the Instrumentl "Matches" tab).  If the API does not support
        this parameter the call still succeeds and returns the first page of
        all available grants.
        """
        params = {"page_size": 50}
        if project_id:
            params["project_id"] = project_id
        result = self._make_request("/v1/grants", params)
        if not result:
            return []
        return result.get("grants", [])

    def get_grant(self, grant_id):
        return self._make_request(f"/v1/grants/{grant_id}")

    def get_saved_grants(self, page_size=50, cursor=None, project_id=None):
        params = {"page_size": page_size}
        if cursor:
            params["cursor"] = cursor
        if project_id:
            params["project_id"] = project_id
        return self._make_request("/v1/saved_grants", params)

    def get_funders(self, page_size=50, cursor=None, search=None):
        params = {"page_size": page_size}
        if cursor:
            params["cursor"] = cursor
        if search:
            params["search"] = search
        return self._make_request("/v1/funders", params)

    def get_funder(self, funder_id):
        return self._make_request(f"/v1/funders/{funder_id}")

    def get_all_grants(self, callback=None):
        all_grants = []
        seen_ids = set()
        cursor = None
        page = 1
        while True:
            if callback:
                callback(f"Fetching grants page {page}...")
            result = self.get_grants(page_size=50, cursor=cursor)
            if not result:
                break
            grants = result.get('grants', [])
            for g in grants:
                gid = g.get('id')
                if gid not in seen_ids:
                    seen_ids.add(gid)
                    all_grants.append(g)
            meta = result.get('meta', {})
            if not meta.get('has_more', False):
                break
            cursor = meta.get('cursor')
            page += 1
            time.sleep(0.25)
        return all_grants

    def get_all_saved_grants(self, project_id=None, callback=None):
        all_saved = []
        seen_ids = set()
        cursor = None
        page = 1
        while True:
            if callback:
                callback(f"Fetching saved grants page {page}...")
            result = self.get_saved_grants(page_size=50, cursor=cursor, project_id=project_id)
            if not result:
                break
            saved = result.get('saved_grants', [])
            for s in saved:
                sid = s.get('id') or s.get('grant_id')
                if sid not in seen_ids:
                    seen_ids.add(sid)
                    all_saved.append(s)
            meta = result.get('meta', {})
            if not meta.get('has_more', False):
                break
            cursor = meta.get('cursor')
            page += 1
            time.sleep(0.25)
        return all_saved


# ==============================================================================
# GUI APPLICATION
# ==============================================================================

class GrantMatcherApp:
    """Main application GUI."""

    COLORS = {
        'primary': '#6366f1',
        'primary_hover': '#4f46e5',
        'primary_light': '#e0e7ff',
        'secondary': '#64748b',
        'success': '#10b981',
        'warning': '#f59e0b',
        'danger': '#ef4444',
        'background': '#f8fafc',
        'surface': '#ffffff',
        'text': '#1e293b',
        'text_secondary': '#64748b',
        'border': '#e2e8f0',
        'accent': '#8b5cf6',
    }

    def __init__(self, root):
        self.root = root
        self.root.title("Instrumentl Grant Matcher v2.0")
        self.root.geometry("950x750")
        self.root.minsize(850, 650)
        self.root.configure(bg=self.COLORS['background'])

        self.config = load_config()
        self.uploaded_files = []
        self.grants_data = []
        self.match_results = []
        self.api_client = None
        self.projects_list = []
        self.selected_project_id = None

        self.setup_modern_styles()
        self.create_wizard()

    def setup_modern_styles(self):
        style = ttk.Style()
        available_themes = style.theme_names()
        if 'clam' in available_themes:
            style.theme_use('clam')

        style.configure('.', font=('Segoe UI', 10), background=self.COLORS['background'])
        style.configure('TFrame', background=self.COLORS['background'])
        style.configure('Card.TFrame', background=self.COLORS['surface'], relief='flat')
        style.configure('TLabel', background=self.COLORS['background'], foreground=self.COLORS['text'],
                        font=('Segoe UI', 10))
        style.configure('Title.TLabel', font=('Segoe UI', 24, 'bold'), foreground=self.COLORS['primary'],
                        background=self.COLORS['background'])
        style.configure('Subtitle.TLabel', font=('Segoe UI', 11), foreground=self.COLORS['text_secondary'],
                        background=self.COLORS['background'])
        style.configure('Header.TLabel', font=('Segoe UI', 12, 'bold'), foreground=self.COLORS['text'],
                        background=self.COLORS['surface'])
        style.configure('Success.TLabel', foreground=self.COLORS['success'], font=('Segoe UI', 10, 'bold'))
        style.configure('Status.TLabel', background=self.COLORS['surface'], foreground=self.COLORS['text_secondary'],
                        font=('Segoe UI', 9), padding=(10, 5))
        style.configure('TButton', font=('Segoe UI', 10), padding=(16, 8), background=self.COLORS['surface'],
                        foreground=self.COLORS['text'])
        style.map('TButton', background=[('active', self.COLORS['border'])],
                  relief=[('pressed', 'flat'), ('!pressed', 'flat')])
        style.configure('Primary.TButton', font=('Segoe UI', 11, 'bold'), padding=(20, 12),
                        background=self.COLORS['primary'], foreground='white')
        style.map('Primary.TButton',
                  background=[('active', self.COLORS['primary_hover']), ('disabled', self.COLORS['border'])],
                  foreground=[('disabled', self.COLORS['text_secondary'])])
        style.configure('Secondary.TButton', font=('Segoe UI', 10), padding=(12, 6), background=self.COLORS['surface'])
        style.configure('TNotebook', background=self.COLORS['background'], borderwidth=0)
        style.configure('TNotebook.Tab', font=('Segoe UI', 10), padding=(20, 10), background=self.COLORS['surface'],
                        foreground=self.COLORS['text_secondary'])
        style.map('TNotebook.Tab', background=[('selected', self.COLORS['primary_light'])],
                  foreground=[('selected', self.COLORS['primary'])], expand=[('selected', [0, 0, 0, 2])])
        style.configure('TEntry', font=('Segoe UI', 10), padding=(8, 6), fieldbackground=self.COLORS['surface'])
        style.configure('TProgressbar', background=self.COLORS['primary'], troughcolor=self.COLORS['border'],
                        borderwidth=0, thickness=8)
        style.configure('TCheckbutton', background=self.COLORS['surface'], foreground=self.COLORS['text'],
                        font=('Segoe UI', 10))
        style.configure('TCombobox', font=('Segoe UI', 10), padding=(8, 6))

    def create_wizard(self):
        self.main_frame = ttk.Frame(self.root, padding="20")
        self.main_frame.pack(fill=tk.BOTH, expand=True)

        # Header
        header_frame = ttk.Frame(self.main_frame)
        header_frame.pack(fill=tk.X, pady=(0, 15))
        title_container = ttk.Frame(header_frame)
        title_container.pack(anchor=tk.W)
        ttk.Label(title_container, text="📊", font=('Segoe UI', 28), background=self.COLORS['background']).pack(
            side=tk.LEFT, padx=(0, 10))
        title_text_frame = ttk.Frame(title_container)
        title_text_frame.pack(side=tk.LEFT)
        ttk.Label(title_text_frame, text="Instrumentl Grant Matcher", style='Title.TLabel').pack(anchor=tk.W)
        ttk.Label(title_text_frame, text="Match your documents with relevant grant opportunities",
                  style='Subtitle.TLabel').pack(anchor=tk.W)

        ttk.Separator(self.main_frame, orient='horizontal').pack(fill=tk.X, pady=(0, 15))

        # Navigation buttons (packed bottom first so they stay visible)
        nav_frame = ttk.Frame(self.main_frame)
        nav_frame.pack(side=tk.BOTTOM, fill=tk.X, pady=(10, 0))

        self.prev_btn = ttk.Button(nav_frame, text="◀ Previous", command=self.previous_tab, style='Secondary.TButton')
        self.prev_btn.pack(side=tk.LEFT, padx=(0, 10))

        self.next_btn = ttk.Button(nav_frame, text="Next ▶", command=self.next_tab, style='Primary.TButton')
        self.next_btn.pack(side=tk.RIGHT)

        # Status bar (packed bottom before notebook)
        status_frame = ttk.Frame(self.main_frame)
        status_frame.pack(side=tk.BOTTOM, fill=tk.X, pady=(10, 0))
        self.status_var = tk.StringVar(value="✓ Ready")
        ttk.Label(status_frame, textvariable=self.status_var, style='Status.TLabel').pack(fill=tk.X)

        # Notebook (fills remaining space between header and bottom buttons)
        self.notebook = ttk.Notebook(self.main_frame)
        self.notebook.pack(fill=tk.BOTH, expand=True, pady=(0, 10))

        self.create_setup_tab()
        self.create_documents_tab()
        self.create_fetch_tab()
        self.create_match_tab()
        self.create_results_tab()

        # Update button states based on current tab
        self.notebook.bind("<<NotebookTabChanged>>", self.update_nav_buttons)
        self.update_nav_buttons()

    def next_tab(self):
        """Navigate to the next tab"""
        current = self.notebook.index(self.notebook.select())
        if current < self.notebook.index("end") - 1:
            self.notebook.select(current + 1)

    def previous_tab(self):
        """Navigate to the previous tab"""
        current = self.notebook.index(self.notebook.select())
        if current > 0:
            self.notebook.select(current - 1)

    def update_nav_buttons(self, event=None):
        """Update Previous/Next button states based on current tab"""
        current = self.notebook.index(self.notebook.select())
        total_tabs = self.notebook.index("end")

        # Disable Previous on first tab
        if current == 0:
            self.prev_btn.state(['disabled'])
        else:
            self.prev_btn.state(['!disabled'])

        # Disable Next on last tab
        if current == total_tabs - 1:
            self.next_btn.state(['disabled'])
        else:
            self.next_btn.state(['!disabled'])

    def _on_canvas_configure(self, event):
        """Update the canvas window width when the canvas is resized"""
        self.canvas.itemconfig(self.canvas_frame, width=event.width)

    def _on_mousewheel(self, event):
        """Handle mousewheel scrolling"""
        if event.num == 5 or event.delta < 0:
            self.canvas.yview_scroll(1, "units")
        elif event.num == 4 or event.delta > 0:
            self.canvas.yview_scroll(-1, "units")

    def create_setup_tab(self):
        outer_frame = ttk.Frame(self.notebook)
        self.notebook.add(outer_frame, text="  1. API Setup  ")
        frame = ttk.Frame(outer_frame, padding="30", style='Card.TFrame')
        frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

        ttk.Label(frame, text="🔐 Connect to Instrumentl", font=('Segoe UI', 14, 'bold'), foreground=self.COLORS['text'],
                  background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 5))
        ttk.Label(frame, text="Enter your API credentials from Instrumentl → Integrations → API", font=('Segoe UI', 10),
                  foreground=self.COLORS['text_secondary'], background=self.COLORS['surface'], wraplength=600).pack(
            anchor=tk.W, pady=(0, 25))

        form_frame = ttk.Frame(frame, style='Card.TFrame')
        form_frame.pack(fill=tk.X)

        ttk.Label(form_frame, text="API Key ID", font=('Segoe UI', 10, 'bold'), foreground=self.COLORS['text'],
                  background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 5))
        self.api_key_id_var = tk.StringVar(value=self.config.get('api_key_id', ''))
        ttk.Entry(form_frame, textvariable=self.api_key_id_var, width=65, font=('Consolas', 10)).pack(fill=tk.X,
                                                                                                      pady=(0, 15))

        ttk.Label(form_frame, text="Private Key", font=('Segoe UI', 10, 'bold'), foreground=self.COLORS['text'],
                  background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 5))
        self.api_private_key_var = tk.StringVar(value=self.config.get('api_private_key', ''))
        ttk.Entry(form_frame, textvariable=self.api_private_key_var, width=65, show='•', font=('Consolas', 10)).pack(
            fill=tk.X, pady=(0, 20))

        btn_frame = ttk.Frame(frame, style='Card.TFrame')
        btn_frame.pack(anchor=tk.W, pady=(10, 0))
        ttk.Button(btn_frame, text="🔗 Test Connection", command=self.test_connection, style='Primary.TButton').pack(
            side=tk.LEFT, padx=(0, 10))
        ttk.Button(btn_frame, text="💾 Save Credentials", command=self.save_credentials, style='Secondary.TButton').pack(
            side=tk.LEFT)

        self.connection_status_var = tk.StringVar(value="")
        ttk.Label(frame, textvariable=self.connection_status_var, font=('Segoe UI', 10),
                  background=self.COLORS['surface'], wraplength=500).pack(anchor=tk.W, pady=(20, 0))

    def create_documents_tab(self):
        outer_frame = ttk.Frame(self.notebook)
        self.notebook.add(outer_frame, text="  2. Documents  ")
        frame = ttk.Frame(outer_frame, padding="20", style='Card.TFrame')
        frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

        ttk.Label(frame, text="📄 Upload Your Documents", font=('Segoe UI', 14, 'bold'), foreground=self.COLORS['text'],
                  background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 5))
        ttk.Label(frame,
                  text="Add documents that describe your organization, programs, and funding needs.\nSupported: PDF, Word, Excel, PowerPoint, CSV, and text files",
                  font=('Segoe UI', 10), foreground=self.COLORS['text_secondary'], background=self.COLORS['surface'],
                  wraplength=700).pack(anchor=tk.W, pady=(0, 10))

        list_frame = ttk.Frame(frame, style='Card.TFrame')
        list_frame.pack(fill=tk.X, pady=(0, 10))
        listbox_container = tk.Frame(list_frame, bg=self.COLORS['border'], padx=1, pady=1)
        listbox_container.pack(fill=tk.X)

        self.file_listbox = tk.Listbox(listbox_container, height=4, selectmode=tk.EXTENDED, font=('Segoe UI', 10),
                                       bg=self.COLORS['surface'], fg=self.COLORS['text'],
                                       selectbackground=self.COLORS['primary_light'],
                                       selectforeground=self.COLORS['primary'], borderwidth=0, highlightthickness=0,
                                       activestyle='none')
        scrollbar = ttk.Scrollbar(listbox_container, orient=tk.VERTICAL, command=self.file_listbox.yview)
        self.file_listbox.configure(yscrollcommand=scrollbar.set)
        self.file_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=1, pady=1)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        btn_frame = ttk.Frame(frame, style='Card.TFrame')
        btn_frame.pack(fill=tk.X, pady=(0, 10))
        ttk.Button(btn_frame, text="📁 Add Files", command=self.add_files, style='Primary.TButton').pack(side=tk.LEFT,
                                                                                                        padx=(0, 8))
        ttk.Button(btn_frame, text="📂 Add Folder", command=self.add_folder, style='Secondary.TButton').pack(
            side=tk.LEFT, padx=(0, 8))
        ttk.Button(btn_frame, text="✕ Remove", command=self.remove_files, style='Secondary.TButton').pack(side=tk.LEFT,
                                                                                                          padx=(0, 8))
        ttk.Button(btn_frame, text="🗑 Clear All", command=self.clear_files, style='Secondary.TButton').pack(
            side=tk.LEFT)

        settings_frame = ttk.Frame(frame, style='Card.TFrame')
        settings_frame.pack(fill=tk.X, pady=(0, 0))
        ttk.Label(settings_frame, text="Chunk Size:", font=('Segoe UI', 10, 'bold'), foreground=self.COLORS['text'],
                  background=self.COLORS['surface']).pack(side=tk.LEFT)
        self.chunk_size_var = tk.StringVar(value=str(self.config.get('chunk_size', 500)))
        ttk.Entry(settings_frame, textvariable=self.chunk_size_var, width=8, font=('Segoe UI', 10)).pack(side=tk.LEFT,
                                                                                                         padx=(8, 8))
        ttk.Label(settings_frame, text="words per paragraph", font=('Segoe UI', 9),
                  foreground=self.COLORS['text_secondary'], background=self.COLORS['surface']).pack(side=tk.LEFT)

    def create_fetch_tab(self):
        outer_frame = ttk.Frame(self.notebook)
        self.notebook.add(outer_frame, text="  3. Fetch Grants  ")
        frame = ttk.Frame(outer_frame, padding="30", style='Card.TFrame')
        frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

        ttk.Label(frame, text="☁️ Fetch from Instrumentl", font=('Segoe UI', 14, 'bold'),
                  foreground=self.COLORS['text'], background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 5))
        ttk.Label(frame, text="Download grant opportunities from your Instrumentl account", font=('Segoe UI', 10),
                  foreground=self.COLORS['text_secondary'], background=self.COLORS['surface'], wraplength=600).pack(
            anchor=tk.W, pady=(0, 25))

        # Project selection
        project_container = tk.Frame(frame, bg=self.COLORS['border'], padx=1, pady=1)
        project_container.pack(fill=tk.X, pady=(0, 15))
        project_frame = ttk.Frame(project_container, padding="15", style='Card.TFrame')
        project_frame.pack(fill=tk.X)
        ttk.Label(project_frame, text="Select Project (Optional)", font=('Segoe UI', 11, 'bold'),
                  foreground=self.COLORS['text'], background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 10))

        project_select_frame = ttk.Frame(project_frame, style='Card.TFrame')
        project_select_frame.pack(fill=tk.X)
        self.project_combo_var = tk.StringVar(value="-- All Projects --")
        self.project_combo = ttk.Combobox(project_select_frame, textvariable=self.project_combo_var, state='readonly',
                                          width=50, font=('Segoe UI', 10))
        self.project_combo['values'] = ["-- All Projects --"]
        self.project_combo.pack(side=tk.LEFT, padx=(0, 10))
        ttk.Button(project_select_frame, text="🔄 Refresh Projects", command=self.refresh_projects,
                   style='Secondary.TButton').pack(side=tk.LEFT)

        # Options
        options_container = tk.Frame(frame, bg=self.COLORS['border'], padx=1, pady=1)
        options_container.pack(fill=tk.X, pady=(0, 20))
        options_frame = ttk.Frame(options_container, padding="20", style='Card.TFrame')
        options_frame.pack(fill=tk.X)
        ttk.Label(options_frame, text="Fetch Options", font=('Segoe UI', 11, 'bold'), foreground=self.COLORS['text'],
                  background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 10))

        self.fetch_saved_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(options_frame, text="Fetch Saved Grants (grants you've already saved to projects)",
                        variable=self.fetch_saved_var, style='TCheckbutton').pack(anchor=tk.W, pady=3)
        ttk.Label(options_frame,
                  text="💡 To include Instrumentl's recommended matches, use the auto-save script to save them to your project first.",
                  font=('Segoe UI', 8), foreground=self.COLORS['text_secondary'],
                  background=self.COLORS['surface'], wraplength=650).pack(anchor=tk.W, pady=(2, 0))

        # Location filter
        ttk.Separator(options_frame, orient='horizontal').pack(fill=tk.X, pady=(15, 10))
        ttk.Label(options_frame, text="Geographic Filter", font=('Segoe UI', 10, 'bold'),
                  foreground=self.COLORS['text'], background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 8))

        location_frame = ttk.Frame(options_frame, style='Card.TFrame')
        location_frame.pack(fill=tk.X, pady=(0, 5))

        self.location_filter_var = tk.StringVar(value="all")
        ttk.Radiobutton(location_frame, text="All locations", variable=self.location_filter_var,
                        value="all", style='TCheckbutton').pack(anchor=tk.W, pady=2)
        ttk.Radiobutton(location_frame, text="Indiana only", variable=self.location_filter_var,
                        value="indiana", style='TCheckbutton').pack(anchor=tk.W, pady=2)
        ttk.Radiobutton(location_frame, text="USA (nationwide)", variable=self.location_filter_var,
                        value="usa", style='TCheckbutton').pack(anchor=tk.W, pady=2)
        ttk.Radiobutton(location_frame, text="Indiana + USA nationwide", variable=self.location_filter_var,
                        value="indiana_usa", style='TCheckbutton').pack(anchor=tk.W, pady=2)

        ttk.Label(options_frame,
                  text="Note: Location filtering happens after fetching, based on grant geographic restrictions.",
                  font=('Segoe UI', 8), foreground=self.COLORS['text_secondary'],
                  background=self.COLORS['surface'], wraplength=650).pack(anchor=tk.W, pady=(5, 0))

        ttk.Button(frame, text="⬇️ Fetch Grants from Instrumentl", command=self.fetch_grants,
                   style='Primary.TButton').pack(pady=(10, 20))

        progress_frame = ttk.Frame(frame, style='Card.TFrame')
        progress_frame.pack(fill=tk.X, pady=(0, 10))
        self.fetch_progress_var = tk.StringVar(value="")
        ttk.Label(progress_frame, textvariable=self.fetch_progress_var, font=('Segoe UI', 10),
                  foreground=self.COLORS['text_secondary'], background=self.COLORS['surface']).pack(anchor=tk.W,
                                                                                                    pady=(0, 8))
        self.fetch_progress = ttk.Progressbar(progress_frame, mode='indeterminate', length=500)
        self.fetch_progress.pack(fill=tk.X)

        summary_frame = ttk.Frame(frame, style='Card.TFrame')
        summary_frame.pack(fill=tk.X, pady=(20, 0))
        self.grants_summary_var = tk.StringVar(value="📋 No grants loaded yet")
        ttk.Label(summary_frame, textvariable=self.grants_summary_var, font=('Segoe UI', 12, 'bold'),
                  foreground=self.COLORS['primary'], background=self.COLORS['surface']).pack(anchor=tk.W)

    def create_match_tab(self):
        outer_frame = ttk.Frame(self.notebook)
        self.notebook.add(outer_frame, text="  4. Run Matching  ")
        frame = ttk.Frame(outer_frame, padding="30", style='Card.TFrame')
        frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

        ttk.Label(frame, text="🎯 Find Matching Grants", font=('Segoe UI', 14, 'bold'), foreground=self.COLORS['text'],
                  background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 5))
        ttk.Label(frame,
                  text="Run the matching algorithm to find grants relevant to your documents.\nAll processing is done locally using TF-IDF text similarity.\n\nTip: Set Maximum Results to 0 to get ALL matches for comprehensive analysis.",
                  font=('Segoe UI', 10), foreground=self.COLORS['text_secondary'], background=self.COLORS['surface'],
                  wraplength=700).pack(anchor=tk.W, pady=(0, 25))

        settings_container = tk.Frame(frame, bg=self.COLORS['border'], padx=1, pady=1)
        settings_container.pack(fill=tk.X, pady=(0, 25))
        settings_frame = ttk.Frame(settings_container, padding="20", style='Card.TFrame')
        settings_frame.pack(fill=tk.X)
        ttk.Label(settings_frame, text="Match Settings", font=('Segoe UI', 11, 'bold'), foreground=self.COLORS['text'],
                  background=self.COLORS['surface']).pack(anchor=tk.W, pady=(0, 15))

        score_frame = ttk.Frame(settings_frame, style='Card.TFrame')
        score_frame.pack(fill=tk.X, pady=5)
        ttk.Label(score_frame, text="Minimum Match Score (0-1):", font=('Segoe UI', 10),
                  background=self.COLORS['surface']).pack(side=tk.LEFT)
        self.min_score_var = tk.StringVar(value=str(self.config.get('min_match_score', 0.1)))
        ttk.Entry(score_frame, textvariable=self.min_score_var, width=10, font=('Segoe UI', 10)).pack(side=tk.LEFT,
                                                                                                      padx=(10, 10))
        ttk.Label(score_frame, text="(lower = more results)", font=('Segoe UI', 9),
                  foreground=self.COLORS['text_secondary'], background=self.COLORS['surface']).pack(side=tk.LEFT)

        top_frame = ttk.Frame(settings_frame, style='Card.TFrame')
        top_frame.pack(fill=tk.X, pady=5)
        ttk.Label(top_frame, text="Maximum Results:", font=('Segoe UI', 10), background=self.COLORS['surface']).pack(
            side=tk.LEFT)
        self.top_matches_var = tk.StringVar(value=str(self.config.get('top_matches', 100)))
        ttk.Entry(top_frame, textvariable=self.top_matches_var, width=10, font=('Segoe UI', 10)).pack(side=tk.LEFT,
                                                                                                      padx=(10, 10))
        ttk.Label(top_frame, text="(or enter 0 for ALL matches)", font=('Segoe UI', 9),
                  foreground=self.COLORS['text_secondary'], background=self.COLORS['surface']).pack(side=tk.LEFT)

        run_btn_frame = ttk.Frame(frame, style='Card.TFrame')
        run_btn_frame.pack(pady=20)
        run_btn = tk.Button(run_btn_frame, text="▶  RUN MATCHING", command=self.run_matching,
                            font=('Segoe UI', 14, 'bold'), bg=self.COLORS['primary'], fg='white',
                            activebackground=self.COLORS['primary_hover'], activeforeground='white', relief='flat',
                            padx=40, pady=15, cursor='hand2')
        run_btn.pack()
        run_btn.bind('<Enter>', lambda e: run_btn.configure(bg=self.COLORS['primary_hover']))
        run_btn.bind('<Leave>', lambda e: run_btn.configure(bg=self.COLORS['primary']))

        progress_frame = ttk.Frame(frame, style='Card.TFrame')
        progress_frame.pack(fill=tk.X, pady=(10, 0))
        self.match_progress_var = tk.StringVar(value="")
        ttk.Label(progress_frame, textvariable=self.match_progress_var, font=('Segoe UI', 10),
                  foreground=self.COLORS['text_secondary'], background=self.COLORS['surface']).pack(anchor=tk.W,
                                                                                                    pady=(0, 8))
        self.match_progress = ttk.Progressbar(progress_frame, mode='determinate', length=500)
        self.match_progress.pack(fill=tk.X)

    def create_results_tab(self):
        outer_frame = ttk.Frame(self.notebook)
        self.notebook.add(outer_frame, text="  5. Results  ")
        frame = ttk.Frame(outer_frame, padding="30", style='Card.TFrame')
        frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=20)

        header_frame = ttk.Frame(frame, style='Card.TFrame')
        header_frame.pack(fill=tk.X, pady=(0, 15))
        self.results_count_var = tk.StringVar(value="📊 No results yet")
        ttk.Label(header_frame, textvariable=self.results_count_var, font=('Segoe UI', 14, 'bold'),
                  foreground=self.COLORS['text'], background=self.COLORS['surface']).pack(side=tk.LEFT)

        export_frame = ttk.Frame(header_frame, style='Card.TFrame')
        export_frame.pack(side=tk.RIGHT)
        ttk.Button(export_frame, text="📥 Export Excel", command=self.export_excel, style='Primary.TButton').pack(
            side=tk.RIGHT, padx=(8, 0))
        ttk.Button(export_frame, text="📄 Export CSV", command=self.export_results, style='Secondary.TButton').pack(
            side=tk.RIGHT)

        text_container = tk.Frame(frame, bg=self.COLORS['border'], padx=1, pady=1)
        text_container.pack(fill=tk.BOTH, expand=True)
        self.results_text = scrolledtext.ScrolledText(text_container, wrap=tk.WORD, font=('Consolas', 10),
                                                      bg=self.COLORS['surface'], fg=self.COLORS['text'], borderwidth=0,
                                                      highlightthickness=0, padx=15, pady=15)
        self.results_text.pack(fill=tk.BOTH, expand=True)

        self.file_location_var = tk.StringVar(value="")
        ttk.Label(frame, textvariable=self.file_location_var, font=('Segoe UI', 9), foreground=self.COLORS['success'],
                  background=self.COLORS['surface'], wraplength=700).pack(anchor=tk.W, pady=(10, 0))

    # ==================== Event Handlers ====================

    def test_connection(self):
        api_key_id = self.api_key_id_var.get().strip()
        api_private_key = self.api_private_key_var.get().strip()
        if not api_key_id or not api_private_key:
            messagebox.showerror("Error", "Please enter both API Key ID and Private Key")
            return
        self.connection_status_var.set("Testing connection...")
        self.root.update()
        try:
            client = InstrumentlAPI(
                api_key_id, api_private_key,
                max_retries=self.config.get('max_retries', 3),
                retry_base_delay=self.config.get('retry_base_delay', 1.0),
                retry_max_delay=self.config.get('retry_max_delay', 30.0),
            )
            account = client.get_account()
            org_name = account.get('organization_name', 'Unknown')
            self.connection_status_var.set(f"✓ Connected successfully! Organization: {org_name}")
            self.api_client = client
            self.refresh_projects()
        except Exception as e:
            self.connection_status_var.set(f"✗ Connection failed: {str(e)}")
            messagebox.showerror("Connection Error", str(e))

    def save_credentials(self):
        self.config['api_key_id'] = self.api_key_id_var.get().strip()
        self.config['api_private_key'] = self.api_private_key_var.get().strip()
        save_config(self.config)
        messagebox.showinfo("Saved", "Credentials saved successfully!")

    def refresh_projects(self):
        api_key_id = self.api_key_id_var.get().strip()
        api_private_key = self.api_private_key_var.get().strip()
        if not api_key_id or not api_private_key:
            messagebox.showerror("Error", "Please configure API credentials first")
            return

        self.status_var.set("Refreshing projects...")
        self.root.update()

        try:
            client = InstrumentlAPI(
                api_key_id, api_private_key,
                max_retries=self.config.get('max_retries', 3),
                retry_base_delay=self.config.get('retry_base_delay', 1.0),
                retry_max_delay=self.config.get('retry_max_delay', 30.0),
            )
            projects = client.get_all_projects()

            self.projects_list = projects
            project_names = ["-- All Projects --"]

            for p in projects:
                # API uses 'project_title' field according to documentation
                name = p.get('project_title') or f"Project {p.get('id', 'Unknown')}"
                if name and name.strip():  # Only add if not empty
                    project_names.append(name)

            self.project_combo['values'] = project_names
            self.project_combo_var.set("-- All Projects --")
            loaded_count = len(project_names) - 1  # subtract 1 for "-- All Projects --"
            self.status_var.set(f"✓ Loaded {loaded_count} projects")
            messagebox.showinfo("Projects Loaded", "Projects loaded successfully")
        except Exception as e:
            self.status_var.set("✗ Error loading projects")
            messagebox.showerror("Error", f"Failed to fetch projects: {str(e)}")

    def get_selected_project_id(self):
        """Get the project ID from the selected project name"""
        selection = self.project_combo_var.get()
        if selection == "-- All Projects --" or not selection:
            return None

        # Find the project by matching the project_title
        for project in self.projects_list:
            # API uses 'project_title' field
            name = project.get('project_title') or f"Project {project.get('id', 'Unknown')}"
            if name and name.strip() == selection:
                return project.get('id')

        return None

    def add_files(self):
        filetypes = [
            ("All Supported", "*.pdf *.docx *.doc *.xlsx *.xls *.pptx *.ppt *.csv *.txt *.md"),
            ("PDF", "*.pdf"), ("Word", "*.docx *.doc"), ("Excel", "*.xlsx *.xls"),
            ("PowerPoint", "*.pptx *.ppt"), ("CSV", "*.csv"), ("Text", "*.txt *.md"), ("All Files", "*.*")
        ]
        files = filedialog.askopenfilenames(filetypes=filetypes)
        for f in files:
            if f not in self.uploaded_files:
                self.uploaded_files.append(f)
                self.file_listbox.insert(tk.END, os.path.basename(f))

    def add_folder(self):
        folder = filedialog.askdirectory()
        if folder:
            count = 0
            for root, dirs, files in os.walk(folder):
                for f in files:
                    ext = Path(f).suffix.lower()
                    if ext in DocumentProcessor.SUPPORTED_EXTENSIONS:
                        filepath = os.path.join(root, f)
                        if filepath not in self.uploaded_files:
                            self.uploaded_files.append(filepath)
                            self.file_listbox.insert(tk.END, f)
                            count += 1
            messagebox.showinfo("Added", f"Added {count} files from folder")

    def remove_files(self):
        selected = self.file_listbox.curselection()
        for idx in reversed(selected):
            self.file_listbox.delete(idx)
            del self.uploaded_files[idx]

    def clear_files(self):
        self.file_listbox.delete(0, tk.END)
        self.uploaded_files.clear()

    def fetch_grants(self):
        api_key_id = self.api_key_id_var.get().strip()
        api_private_key = self.api_private_key_var.get().strip()
        if not api_key_id or not api_private_key:
            messagebox.showerror("Error", "Please configure API credentials first")
            self.notebook.select(0)
            return

        if not self.fetch_saved_var.get():
            messagebox.showerror("Error", "Please select at least one fetch option")
            return

        selected_project_id = self.get_selected_project_id()
        location_filter = self.location_filter_var.get()

        def fetch_thread():
            try:
                self.fetch_progress.start()
                self.status_var.set("⏳ Fetching grants...")
                client = InstrumentlAPI(
                    api_key_id, api_private_key,
                    max_retries=self.config.get('max_retries', 3),
                    retry_base_delay=self.config.get('retry_base_delay', 1.0),
                    retry_max_delay=self.config.get('retry_max_delay', 30.0),
                )
                all_grants = []

                if self.fetch_saved_var.get():
                    self.fetch_progress_var.set("Fetching saved grants...")
                    self.root.update()
                    saved = client.get_all_saved_grants(
                        project_id=selected_project_id,
                        callback=lambda msg: self.update_fetch_status(msg)
                    )
                    for s in saved:
                        grant_id = s.get('grant_id')
                        if grant_id:
                            try:
                                grant_detail = client.get_grant(grant_id)
                                if grant_detail:
                                    grant_detail['_saved_grant_info'] = s
                                    all_grants.append(grant_detail)
                                time.sleep(0.2)
                            except:
                                pass

                # Apply location filter
                if location_filter != "all":
                    self.fetch_progress_var.set("Applying location filter...")
                    self.root.update()
                    filtered_grants = []
                    for grant in all_grants:
                        if self.grant_matches_location(grant, location_filter):
                            filtered_grants.append(grant)

                    filtered_count = len(all_grants) - len(filtered_grants)
                    all_grants = filtered_grants
                    filter_msg = f" ({filtered_count} grants filtered by location)" if filtered_count > 0 else ""
                else:
                    filter_msg = ""

                self.grants_data = all_grants
                self.fetch_progress.stop()
                self.fetch_progress_var.set("✓ Complete!")
                project_info = f" from '{self.project_combo_var.get()}'" if selected_project_id else ""
                self.grants_summary_var.set(f"📊 Loaded {len(all_grants)} grants{project_info}{filter_msg}")
                self.status_var.set(f"✓ Successfully loaded {len(all_grants)} grants{filter_msg}")

                if len(all_grants) == 0:
                    messagebox.showinfo("No Grants Found",
                                        "No grants were found.\n\nTips:\n• Make sure you have saved grants in Instrumentl\n• Try selecting 'All Projects'\n• Check both fetch options\n• Try changing the location filter")
                else:
                    messagebox.showinfo("Success", f"Successfully loaded {len(all_grants)} grants!{filter_msg}")

            except Exception as e:
                self.fetch_progress.stop()
                self.fetch_progress_var.set(f"✗ Error: {str(e)}")
                self.status_var.set("✗ Error fetching grants")
                messagebox.showerror("Fetch Error", str(e))

        thread = threading.Thread(target=fetch_thread, daemon=True)
        thread.start()

    def grant_matches_location(self, grant, location_filter):
        """Check if grant matches the location filter"""
        # Get location information from grant categories
        categories = grant.get('categories', {})
        if not isinstance(categories, dict):
            # If no location data, include by default
            return True

        # Check for location categories
        geographic_area = categories.get('geographic_area_category', [])
        if not geographic_area:
            # No location restrictions = available everywhere
            return True

        # Convert to lowercase for case-insensitive matching
        locations = [loc.lower() if isinstance(loc, str) else '' for loc in geographic_area]

        if location_filter == "indiana":
            # Only Indiana
            return any('indiana' in loc or 'in' == loc for loc in locations)

        elif location_filter == "usa":
            # USA nationwide (not state-specific)
            # Look for terms like "national", "nationwide", "united states", "usa", "all states"
            nationwide_terms = ['national', 'nationwide', 'united states', 'usa', 'all states', 'u.s.']
            return any(any(term in loc for term in nationwide_terms) for loc in locations)

        elif location_filter == "indiana_usa":
            # Indiana OR USA nationwide
            has_indiana = any('indiana' in loc or 'in' == loc for loc in locations)
            nationwide_terms = ['national', 'nationwide', 'united states', 'usa', 'all states', 'u.s.']
            has_nationwide = any(any(term in loc for term in nationwide_terms) for loc in locations)
            return has_indiana or has_nationwide

        return True

    def update_fetch_status(self, msg):
        """Thread-safe status update"""
        self.fetch_progress_var.set(msg)
        self.root.update_idletasks()

    def run_matching(self):
        if not self.uploaded_files:
            messagebox.showerror("Error", "Please upload documents first")
            self.notebook.select(1)
            return
        if not self.grants_data:
            messagebox.showerror("Error", "Please fetch grants first")
            self.notebook.select(2)
            return

        def match_thread():
            try:
                self.match_progress['value'] = 0
                self.match_progress_var.set("Processing documents...")
                self.status_var.set("⏳ Running matching algorithm...")
                chunk_size = int(self.chunk_size_var.get())
                min_score = float(self.min_score_var.get())
                top_matches = int(self.top_matches_var.get())

                self.config['chunk_size'] = chunk_size
                self.config['min_match_score'] = min_score
                self.config['top_matches'] = top_matches
                save_config(self.config)

                doc_chunks = []
                doc_metadata = []
                total_files = len(self.uploaded_files)

                for idx, filepath in enumerate(self.uploaded_files):
                    self.match_progress_var.set(f"Processing: {os.path.basename(filepath)}")
                    self.match_progress['value'] = (idx / total_files) * 30
                    self.root.update()
                    try:
                        text = DocumentProcessor.extract_text(filepath)
                        chunks = TextChunker.chunk_text(text, chunk_size=chunk_size)
                        for i, chunk in enumerate(chunks):
                            doc_chunks.append(chunk)
                            doc_metadata.append(
                                {'file': os.path.basename(filepath), 'chunk_index': i, 'total_chunks': len(chunks)})
                    except Exception as e:
                        print(f"Error processing {filepath}: {e}")

                if not doc_chunks:
                    messagebox.showerror("Error", "No text could be extracted from documents")
                    self.status_var.set("✗ Error: No text extracted")
                    return

                combined_doc_text = ' '.join(doc_chunks)
                self.match_progress_var.set("Building grant index...")
                self.match_progress['value'] = 40
                self.root.update()

                matcher = TFIDFMatcher()
                grant_texts = []
                grant_metas = []

                for grant in self.grants_data:
                    text_parts = [grant.get('name', ''), grant.get('overview', ''),
                                  grant.get('funder', '') if isinstance(grant.get('funder'), str) else '']
                    funder = grant.get('funder')
                    if isinstance(funder, dict):
                        text_parts.append(funder.get('name', ''))
                    categories = grant.get('categories', {})
                    if isinstance(categories, dict):
                        for cat_type, cat_values in categories.items():
                            if isinstance(cat_values, list):
                                text_parts.extend(cat_values)
                    grant_text = ' '.join(str(p) for p in text_parts if p)
                    if grant_text.strip():
                        grant_texts.append(grant_text)
                        grant_metas.append(grant)

                matcher.add_documents(grant_texts, grant_metas)
                self.match_progress_var.set("Building search index...")
                self.match_progress['value'] = 60
                self.root.update()
                matcher.build_index()

                self.match_progress_var.set("Finding matches...")
                self.match_progress['value'] = 80
                self.root.update()

                # If top_matches is 0, get ALL matches above min_score
                actual_top_k = top_matches
                if top_matches == 0:
                    actual_top_k = len(grant_metas)  # Set to total number of grants
                    self.match_progress_var.set(f"Finding ALL matches (analyzing {len(grant_metas)} grants)...")
                    self.root.update()

                matches = matcher.find_matches(combined_doc_text, top_k=actual_top_k, min_score=min_score)
                self.match_results = matches

                self.match_progress_var.set("✓ Complete!")
                self.match_progress['value'] = 100

                # Detailed status message
                if top_matches == 0:
                    self.status_var.set(
                        f"✓ Found {len(matches)} grants above score {min_score} (analyzed {len(grant_metas)} total grants)")
                else:
                    self.status_var.set(f"✓ Found top {len(matches)} matching grants (out of {len(grant_metas)} total)")

                self.display_results()
                self.notebook.select(4)
            except Exception as e:
                self.match_progress_var.set(f"✗ Error: {str(e)}")
                self.status_var.set("✗ Error during matching")
                messagebox.showerror("Matching Error", str(e))
                import traceback
                traceback.print_exc()

        thread = threading.Thread(target=match_thread, daemon=True)
        thread.start()

    def display_results(self):
        self.results_text.delete(1.0, tk.END)
        self.file_location_var.set("")

        if not self.match_results:
            self.results_count_var.set("📊 No matches found")
            self.results_text.insert(tk.END, "No matching grants found.\n\n")
            self.results_text.insert(tk.END, "💡 Tips:\n")
            self.results_text.insert(tk.END, "  • Lower the minimum match score (try 0.05)\n")
            self.results_text.insert(tk.END, "  • Add more descriptive documents\n")
            self.results_text.insert(tk.END, "  • Include mission statements and program descriptions\n")
            return

        self.results_count_var.set(f"🎉 Found {len(self.match_results)} matching grants")

        for idx, match in enumerate(self.match_results, 1):
            grant = match['metadata']
            score = match['score']
            self.results_text.insert(tk.END, f"\n{'━' * 70}\n")
            self.results_text.insert(tk.END, f"  #{idx}  ")
            score_pct = int(score * 20)
            score_bar = "█" * score_pct + "░" * (20 - score_pct)
            self.results_text.insert(tk.END, f"Match: [{score_bar}] {score:.1%}\n")
            self.results_text.insert(tk.END, f"{'━' * 70}\n\n")
            self.results_text.insert(tk.END, f"  📋 Grant: {grant.get('name', 'N/A')}\n")
            self.results_text.insert(tk.END, f"  🔑 ID: {grant.get('id', 'N/A')}\n")
            funder = grant.get('funder', 'N/A')
            if isinstance(funder, dict):
                funder = funder.get('name', 'N/A')
            self.results_text.insert(tk.END, f"  🏛️ Funder: {funder}\n")
            self.results_text.insert(tk.END, f"  📅 Deadline: {grant.get('next_deadline_date', 'N/A')}\n")
            self.results_text.insert(tk.END, f"  📌 Status: {grant.get('status', 'N/A')}\n")
            overview = grant.get('overview', '')
            if overview:
                if len(overview) > 400:
                    overview = overview[:400] + "..."
                self.results_text.insert(tk.END, f"\n  📝 Overview:\n")
                for line in overview.split('\n'):
                    self.results_text.insert(tk.END, f"     {line}\n")
            self.results_text.insert(tk.END, "\n")

    def export_results(self):
        if not self.match_results:
            messagebox.showwarning("No Results", "No results to export")
            return
        initial_dir = self.config.get('last_export_dir', '')
        if not initial_dir or not os.path.exists(initial_dir):
            initial_dir = os.path.expanduser("~")
        filepath = filedialog.asksaveasfilename(defaultextension=".csv", filetypes=[("CSV Files", "*.csv")],
                                                initialdir=initial_dir,
                                                initialfile=f"grant_matches_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv")
        if filepath:
            try:
                import csv
                with open(filepath, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    writer.writerow(
                        ['Rank', 'Match Score', 'Grant Name', 'Grant ID', 'Funder', 'Next Deadline', 'Status',
                         'Funding Cycle', 'Grant URL', 'Description'])
                    for idx, match in enumerate(self.match_results, 1):
                        grant = match['metadata']
                        funder = grant.get('funder', '')
                        if isinstance(funder, dict):
                            funder = funder.get('name', '')

                        # Get funding cycle interval
                        funding_cycle = ''
                        funding_cycles = grant.get('funding_cycles', [])
                        if funding_cycles and len(funding_cycles) > 0:
                            funding_cycle = funding_cycles[0].get('interval', '')

                        # Get grant URL - construct from slug or ID
                        grant_url = ''
                        slug = grant.get('slug', '')
                        if slug:
                            grant_url = f"https://www.instrumentl.com/grants/{slug}"
                        elif grant.get('id'):
                            # Fallback to ID-based URL
                            grant_url = f"https://www.instrumentl.com/grants/{grant.get('id')}"

                        # Get description/overview
                        description = grant.get('overview', '')
                        if description and len(description) > 1000:
                            description = description[:1000] + '...'

                        writer.writerow(
                            [idx, f"{match['score']:.4f}", grant.get('name', ''), grant.get('id', ''), funder,
                             grant.get('next_deadline_date', ''), grant.get('status', ''),
                             funding_cycle, grant_url, description])
                self.config['last_export_dir'] = os.path.dirname(filepath)
                save_config(self.config)
                self.file_location_var.set(f"✓ File saved to: {filepath}")
                messagebox.showinfo("Exported", f"Results exported successfully!\n\nFile saved to:\n{filepath}")
            except Exception as e:
                messagebox.showerror("Export Error", str(e))

    def export_excel(self):
        if not self.match_results:
            messagebox.showwarning("No Results", "No results to export")
            return
        initial_dir = self.config.get('last_export_dir', '')
        if not initial_dir or not os.path.exists(initial_dir):
            initial_dir = os.path.expanduser("~")
        filepath = filedialog.asksaveasfilename(defaultextension=".xlsx", filetypes=[("Excel Files", "*.xlsx")],
                                                initialdir=initial_dir,
                                                initialfile=f"grant_matches_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx")
        if filepath:
            try:
                import pandas as pd
                data = []
                for idx, match in enumerate(self.match_results, 1):
                    grant = match['metadata']
                    funder = grant.get('funder', '')
                    if isinstance(funder, dict):
                        funder = funder.get('name', '')

                    # Get funding cycle interval
                    funding_cycle = ''
                    funding_cycles = grant.get('funding_cycles', [])
                    if funding_cycles and len(funding_cycles) > 0:
                        funding_cycle = funding_cycles[0].get('interval', '')

                    # Get grant URL - construct from slug or ID
                    grant_url = ''
                    slug = grant.get('slug', '')
                    if slug:
                        grant_url = f"https://www.instrumentl.com/grants/{slug}"
                    elif grant.get('id'):
                        # Fallback to ID-based URL
                        grant_url = f"https://www.instrumentl.com/grants/{grant.get('id')}"

                    data.append({
                        'Rank': idx,
                        'Match Score': match['score'],
                        'Grant Name': grant.get('name', ''),
                        'Grant ID': grant.get('id', ''),
                        'Funder': funder,
                        'Next Deadline': grant.get('next_deadline_date', ''),
                        'Status': grant.get('status', ''),
                        'Is Custom': grant.get('is_custom', False),
                        'Rolling': grant.get('rolling', False),
                        'Funding Cycle': funding_cycle,
                        'Grant URL': grant_url,
                        'Description': grant.get('overview', '')
                    })
                df = pd.DataFrame(data)
                df.to_excel(filepath, index=False, sheet_name='Grant Matches')
                self.config['last_export_dir'] = os.path.dirname(filepath)
                save_config(self.config)
                self.file_location_var.set(f"✓ File saved to: {filepath}")
                messagebox.showinfo("Exported", f"Results exported successfully!\n\nFile saved to:\n{filepath}")
            except ImportError:
                messagebox.showerror("Error", "pandas is required for Excel export. Using CSV instead.")
                self.export_results()
            except Exception as e:
                messagebox.showerror("Export Error", str(e))


# ==============================================================================
# MAIN ENTRY POINT
# ==============================================================================

def main():
    root = tk.Tk()
    try:
        from ctypes import windll
        windll.shcore.SetProcessDpiAwareness(1)
    except:
        pass
    app = GrantMatcherApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
